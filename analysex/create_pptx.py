import os
import pathlib
from pprint import pprint

from pptx import Presentation
from pptx.util import Inches, Cm


def main():
    files = []
    for f in pathlib.Path('all_plots').glob('*.png'):
        if str(f).startswith('all_plots/PIN'):
            files.append(str(f))
    files = sorted(files)

    files_assbly = []
    for f in pathlib.Path('all_plots').glob('*.png'):
        if str(f).startswith('all_plots/ASSBLY'):
            files_assbly.append(str(f))
    files_assbly = sorted(files_assbly)

    files.extend(files_assbly)

    pprint(files)

    prs = Presentation()

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "DETAILED DENSITIES"

    top = Cm(1)
    left = Cm(1)
    height = Inches(6.5)
    for file in files:
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(str(file), left, top, height=height)

    prs.save('densities.pptx')


if __name__ == '__main__':
    os.chdir(pathlib.Path(__file__).resolve().parent.parent)
    main()
